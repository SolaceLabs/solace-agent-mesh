"""
BM25 Document Indexer - Creates individual indexes for each document

This module demonstrates how to:
1. Convert files to text with position tracking:
   - Binary files (PDF, DOCX, PPTX) using MarkItDown
   - Text files (TXT, MD, CSV, HTML, JSON, XML) with line number tracking
2. Chunk documents properly for indexing
3. Create separate BM25 indexes for each document
4. Store metadata for retrieval and citation (including page/slide/line numbers)
5. Support storage-agnostic backends (filesystem, S3, etc.)
"""

import os
import logging
from pathlib import Path
from typing import List, Dict, Any, Tuple, Optional
from abc import ABC, abstractmethod
import json
import pickle
import io
import bm25s
import Stemmer
from markitdown import MarkItDown
# pdfminer installed from markitdown[all]
from pdfminer.high_level import extract_pages
from pdfminer.layout import LTTextContainer

# python-docx installed from markitdown[all]
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not available, Word page tracking will be disabled")

# python-pptx installed from markitdown[all]
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available, PowerPoint slide tracking will be disabled")

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Supported file extensions with their MIME types
SUPPORTED_BINARY_DICT = {
    '.pdf': 'application/pdf',
    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    '.doc': 'application/msword',
    '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    '.ppt': 'application/vnd.ms-powerpoint'
}
SUPPORTED_TEXT_DICT = {
    '.txt': 'text/plain',
    '.md': 'text/markdown',
    '.csv': 'text/csv',
    '.html': 'text/html',
    '.htm': 'text/html',
    '.json': 'application/json',
    '.xml': 'application/xml'
}
supported_binary_extensions_list = list(SUPPORTED_BINARY_DICT.keys())
supported_binary_mime_types_list = list(SUPPORTED_BINARY_DICT.values())
supported_text_extensions_list = list(SUPPORTED_TEXT_DICT.keys())
supported_text_mime_types_list = list(SUPPORTED_TEXT_DICT.values())

class DocumentChunker:
    """
    Handles intelligent document chunking for BM25 indexing.
    
    Chunking Strategy:
    - Split documents into semantic chunks (paragraphs or sections)
    - Maintain overlap between chunks for context continuity
    - Preserve metadata for citation and source tracking
    """
    
    def __init__(
        self,
        chunk_size: int = 512,
        chunk_overlap: int = 128,
        min_chunk_size: int = 50
    ):
        """
        Initialize the document chunker.
        
        Args:
            chunk_size: Target size for each chunk (in characters)
            chunk_overlap: Number of characters to overlap between chunks
            min_chunk_size: Minimum chunk size to keep (avoid tiny chunks)
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.min_chunk_size = min_chunk_size
    
    def chunk_text(
        self, 
        text: str, 
        doc_metadata: Dict[str, Any],
        page_map: Optional[List[Tuple[int, int, int]]] = None,
        position_type: str = 'page'
    ) -> List[Dict[str, Any]]:
        """
        Split text into overlapping chunks with metadata.
        
        Args:
            text: The text content to chunk
            doc_metadata: Metadata about the source document
            page_map: Optional list of (num, char_start, char_end) tuples mapping 
                     character positions to page/slide/line numbers
            position_type: Type of position tracking ('page', 'slide', or 'line')
            
        Returns:
            List of chunk dictionaries with text and metadata
        """
        chunks = []
        
        # Split by paragraphs first (better semantic boundaries)
        paragraphs = text.split('\n\n')
        
        current_chunk = ""
        chunk_index = 0
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            
            # If adding this paragraph exceeds chunk size, save current chunk
            if len(current_chunk) + len(para) > self.chunk_size and current_chunk:
                if len(current_chunk) >= self.min_chunk_size:
                    char_start = len(''.join([c['text'] for c in chunks]))
                    char_end = char_start + len(current_chunk)
                    
                    # Determine position number(s) for this chunk
                    position_numbers = self._get_page_numbers(char_start, char_end, page_map)
                    
                    chunk_data = {
                        'text': current_chunk.strip(),
                        'chunk_index': chunk_index,
                        'doc_name': doc_metadata['filename'],
                        'doc_path': doc_metadata['filepath'],
                        'file_type': doc_metadata['file_type'],
                        'char_start': char_start,
                        'char_end': char_end
                    }
                    
                    # Add position information with appropriate field names
                    if position_numbers:
                        if position_type == 'line':
                            chunk_data['line_numbers'] = position_numbers
                            chunk_data['line_start'] = position_numbers[0]
                            chunk_data['line_end'] = position_numbers[-1]
                        elif position_type == 'slide':
                            chunk_data['slide_numbers'] = position_numbers
                            chunk_data['slide_start'] = position_numbers[0]
                            chunk_data['slide_end'] = position_numbers[-1]
                        else:  # 'page' is default
                            chunk_data['page_numbers'] = position_numbers
                            chunk_data['page_start'] = position_numbers[0]
                            chunk_data['page_end'] = position_numbers[-1]
                    
                    chunks.append(chunk_data)
                    chunk_index += 1
                
                # Start new chunk with overlap
                overlap_text = current_chunk[-self.chunk_overlap:] if len(current_chunk) > self.chunk_overlap else current_chunk
                current_chunk = overlap_text + "\n\n" + para
            else:
                # Add paragraph to current chunk
                if current_chunk:
                    current_chunk += "\n\n" + para
                else:
                    current_chunk = para
        
        # Add the last chunk
        if current_chunk and len(current_chunk) >= self.min_chunk_size:
            char_start = len(''.join([c['text'] for c in chunks]))
            char_end = char_start + len(current_chunk)
            
            # Determine position number(s) for this chunk
            position_numbers = self._get_page_numbers(char_start, char_end, page_map)
            
            chunk_data = {
                'text': current_chunk.strip(),
                'chunk_index': chunk_index,
                'doc_name': doc_metadata['filename'],
                'doc_path': doc_metadata['filepath'],
                'file_type': doc_metadata['file_type'],
                'char_start': char_start,
                'char_end': char_end
            }
            
            # Add position information with appropriate field names
            if position_numbers:
                if position_type == 'line':
                    chunk_data['line_numbers'] = position_numbers
                    chunk_data['line_start'] = position_numbers[0]
                    chunk_data['line_end'] = position_numbers[-1]
                elif position_type == 'slide':
                    chunk_data['slide_numbers'] = position_numbers
                    chunk_data['slide_start'] = position_numbers[0]
                    chunk_data['slide_end'] = position_numbers[-1]
                else:  # 'page' is default
                    chunk_data['page_numbers'] = position_numbers
                    chunk_data['page_start'] = position_numbers[0]
                    chunk_data['page_end'] = position_numbers[-1]
            
            chunks.append(chunk_data)
        
        logger.info(f"Created {len(chunks)} chunks from {doc_metadata['filename']}")
        return chunks
    
    def _get_page_numbers(
        self, 
        char_start: int, 
        char_end: int, 
        page_map: Optional[List[Tuple[int, int, int]]]
    ) -> Optional[List[int]]:
        """
        Determine which page(s) a chunk spans based on character positions.
        
        Args:
            char_start: Starting character position of the chunk
            char_end: Ending character position of the chunk
            page_map: List of (page_num, char_start, char_end) tuples
            
        Returns:
            List of page numbers the chunk spans, or None if no page map available
        """
        if not page_map:
            return None
        
        pages = set()
        for page_num, page_start, page_end in page_map:
            # Check if chunk overlaps with this page
            if not (char_end <= page_start or char_start >= page_end):
                pages.add(page_num)
        
        return sorted(list(pages)) if pages else None


class BM25DocumentIndexer:
    """
    Creates and manages individual BM25 indexes for each document.
    
    Key Features:
    - Converts binary files (PDF, DOCX, PPTX) to text using MarkItDown
    - Tracks page numbers (PDF, DOCX), slide numbers (PPTX), and line numbers (TXT, MD, etc.) for citations
    - Creates separate index for each document (isolated corpus)
    - Uses stemming for better keyword matching
    - Stores metadata for citation and source tracking
    """
    
    def __init__(
        self,
        chunk_size: int = 1024,
        chunk_overlap: int = 128
    ):
        """
        Initialize the BM25 indexer.
        
        Args:
            chunk_size: Size of text chunks for indexing
            chunk_overlap: Overlap between chunks
        """
        
        self.chunker = DocumentChunker(chunk_size, chunk_overlap)
        self.converter = MarkItDown()
        
        # Initialize stemmer for better keyword matching
        self.stemmer = Stemmer.Stemmer("english")
    
    def extract_pdf_with_pages(self, file_path: Path) -> Tuple[str, List[Tuple[int, int, int]], int]:
        """
        Extract text from PDF page by page and create a page mapping.
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            Tuple of (full_text, page_map, total_pages) where:
            - full_text: Combined text from all pages
            - page_map: List of (page_num, char_start, char_end) tuples
            - total_pages: Total number of pages in the PDF
        """
                # Use pdfminer.six for page-by-page extraction
        try:    
            full_text = ""
            page_map = []
            page_num = 0
            
            def extract_text_from_element(element):
                """Recursively extract text from layout elements."""
                text = ""
                if isinstance(element, LTTextContainer):
                    text += element.get_text()
                elif hasattr(element, '__iter__'):
                    for child in element:
                        text += extract_text_from_element(child)
                return text
            
            for page_layout in extract_pages(str(file_path)):
                page_num += 1
                page_text = extract_text_from_element(page_layout)
                
                char_start = len(full_text)
                full_text += page_text
                char_end = len(full_text)
                
                # Add page mapping (page_num, char_start, char_end)
                page_map.append((page_num, char_start, char_end))
                
                # Add separator between pages
                full_text += "\n\n"
            
            logger.info(f"Extracted {page_num} pages from {file_path.name} with page tracking")
            return full_text, page_map, page_num
                
            
        except Exception as e:
            logger.error(f"Error extracting PDF pages from {file_path.name}: {e}")
            return "", [], 0
    
    def extract_docx_with_pages(self, file_path: Path) -> Tuple[str, List[Tuple[int, int, int]], int]:
        """
        Extract text from DOCX file page by page and create a page mapping.
        
        Note: Word documents don't have fixed pages like PDFs. We approximate pages
        based on paragraph breaks and typical page length (around 3000 characters per page).
        
        Args:
            file_path: Path to the DOCX file
            
        Returns:
            Tuple of (full_text, page_map, total_pages) where:
            - full_text: Combined text from all paragraphs
            - page_map: List of (page_num, char_start, char_end) tuples
            - total_pages: Estimated total number of pages
        """
        if not DOCX_AVAILABLE:
            logger.warning(f"python-docx not available, cannot extract pages from {file_path.name}")
            return "", [], 0
        
        try:
            doc = DocxDocument(str(file_path))
            full_text = ""
            page_map = []
            page_num = 1
            
            # Approximate characters per page (typical A4 page with standard formatting)
            # This is an estimation since Word documents are reflowable
            CHARS_PER_PAGE = 3000
            
            current_page_start = 0
            
            # Extract text from all paragraphs
            for paragraph in doc.paragraphs:
                para_text = paragraph.text
                if para_text.strip():  # Only add non-empty paragraphs
                    full_text += para_text + "\n\n"
            
            # Create page mapping based on character count
            total_chars = len(full_text)
            current_pos = 0
            
            while current_pos < total_chars:
                char_start = current_pos
                char_end = min(current_pos + CHARS_PER_PAGE, total_chars)
                
                # Try to break at paragraph boundary (look for \n\n)
                if char_end < total_chars:
                    # Look for the next paragraph break within a reasonable range
                    search_start = max(char_end - 200, char_start)
                    search_end = min(char_end + 200, total_chars)
                    para_break = full_text.find('\n\n', search_start, search_end)
                    if para_break != -1:
                        char_end = para_break + 2  # Include the line breaks
                
                page_map.append((page_num, char_start, char_end))
                page_num += 1
                current_pos = char_end
            
            total_pages = len(page_map)
            
            logger.info(f"Extracted {total_pages} estimated pages from {file_path.name} with page tracking")
            return full_text, page_map, total_pages
            
        except Exception as e:
            logger.error(f"Error extracting DOCX pages from {file_path.name}: {e}")
            return "", [], 0
    
    def extract_pptx_with_slides(self, file_path: Path) -> Tuple[str, List[Tuple[int, int, int]], int]:
        """
        Extract text from PowerPoint file slide by slide and create a slide mapping.
        
        Args:
            file_path: Path to the PPTX/PPT file
            
        Returns:
            Tuple of (full_text, slide_map, total_slides) where:
            - full_text: Combined text from all slides
            - slide_map: List of (slide_num, char_start, char_end) tuples
            - total_slides: Total number of slides in the presentation
        """
        if not PPTX_AVAILABLE:
            logger.warning(f"python-pptx not available, cannot extract slides from {file_path.name}")
            return "", [], 0
        
        try:
            prs = Presentation(str(file_path))
            full_text = ""
            slide_map = []
            slide_num = 0
            
            for slide in prs.slides:
                slide_num += 1
                slide_text = ""
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                    
                    # Also check for text in tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text:
                                    slide_text += cell.text + " "
                            slide_text += "\n"
                
                # Record character positions for this slide
                char_start = len(full_text)
                full_text += slide_text
                char_end = len(full_text)
                
                # Add slide mapping (slide_num, char_start, char_end)
                slide_map.append((slide_num, char_start, char_end))
                
                # Add separator between slides
                full_text += "\n\n"
            
            logger.info(f"Extracted {slide_num} slides from {file_path.name} with slide tracking")
            return full_text, slide_map, slide_num
            
        except Exception as e:
            logger.error(f"Error extracting PPTX slides from {file_path.name}: {e}")
            return "", [], 0
    
    def extract_text_with_lines(self, file_path: Path) -> Tuple[str, List[Tuple[int, int, int]], int]:
        """
        Extract text from a text file and create a line mapping.
        
        Args:
            file_path: Path to the text file
            
        Returns:
            Tuple of (full_text, line_map, total_lines) where:
            - full_text: Combined text from all lines
            - line_map: List of (line_num, char_start, char_end) tuples
            - total_lines: Total number of lines in the file
        """
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                full_text = ""
                line_map = []
                line_num = 0
                
                for line in f:
                    line_num += 1
                    char_start = len(full_text)
                    full_text += line
                    char_end = len(full_text)
                    
                    # Add line mapping (line_num, char_start, char_end)
                    line_map.append((line_num, char_start, char_end))
                
                logger.info(f"Extracted {line_num} lines from {file_path.name} with line tracking")
                return full_text, line_map, line_num
                
        except Exception as e:
            logger.error(f"Error extracting lines from {file_path.name}: {e}")
            return "", [], 0
    
    def convert_to_text(
        self, 
        file_path: Path, 
        supported_type: str,
        file_ext: Optional[str] = None,
        mime_type: Optional[str] = None
    ) -> Tuple[str, str, Optional[List[Tuple[int, int, int]]], Optional[int]]:
        """
        Convert a file to text using MarkItDown.
        
        Args:
            file_path: Path to the file
            file_ext: Optional file extension (used when MIME type is not available)
            mime_type: Optional MIME type of the file (used when file extension is not available)
            supported_type: The supported type of the file ('text' or 'binary')
        Returns:
            Tuple of (text_content, file_type, page_map, total_pages) where:
            - text_content: The extracted text
            - file_type: File extension or MIME type
            - page_map: Optional page/slide/line mapping (list of (num, char_start, char_end))
            - total_pages: Optional total page/slide/line count
        """
        is_pdf = False
        is_docx = False
        is_pptx = False

        # For text files, extract with line tracking
        if supported_type == 'text':
            text_content, line_map, total_lines = self.extract_text_with_lines(file_path)
            if text_content:
                return text_content, file_ext or mime_type, line_map, total_lines
            # Fallback if line extraction fails
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read(), file_ext or mime_type, None, None
        elif supported_type == 'binary':
            is_pdf = (file_ext == '.pdf' or mime_type == 'application/pdf')
            is_docx = (file_ext == '.docx' or mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document')
            is_pptx = (file_ext in ['.pptx', '.ppt'] or mime_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation', 'application/vnd.ms-powerpoint'])
            
            # For PDFs, try to extract with page information
            if is_pdf:
                text_content, page_map, total_pages = self.extract_pdf_with_pages(file_path)
                if text_content:
                    return text_content, file_ext or mime_type, page_map, total_pages
            
            # For DOCX files, try to extract with page information
            if is_docx and DOCX_AVAILABLE:
                text_content, page_map, total_pages = self.extract_docx_with_pages(file_path)
                if text_content:
                    return text_content, file_ext or mime_type, page_map, total_pages
            
            # For PPTX/PPT files, try to extract with slide information
            if is_pptx and PPTX_AVAILABLE:
                text_content, slide_map, total_slides = self.extract_pptx_with_slides(file_path)
                if text_content:
                    return text_content, file_ext or mime_type, slide_map, total_slides
        
        # For binary files (or PDF/DOCX/PPTX fallback), use MarkItDown
        try:
            result = self.converter.convert(str(file_path))
            text_content = result.text_content if result else ""
            logger.info(f"Converted {file_path.name} to text ({len(text_content)} chars)")
            return text_content, file_ext or mime_type, None, None
        except Exception as e:
            logger.error(f"Error converting {file_path.name}: {e}")
            return "", file_ext or mime_type, None, None
        
    def validate_doc_type(self, 
        file_path: Path, 
        mime_type: Optional[str] = None
    ) -> Dict[str, Optional[str]]:
        """
        Validate the document type based on file extension and MIME type.

        Args:
            file_path: Path to the file
            mime_type: Optional MIME type of the file

        Returns:
            Dictionary with supported type, file extension, and MIME type
        """

        file_ext = file_path.suffix.lower()
        mime_type = mime_type.lower() if mime_type else None
        
        supported_type = 'unsupported'
        
        if file_ext in supported_text_extensions_list or mime_type in supported_text_mime_types_list:
            supported_type = 'text'
        elif file_ext in supported_binary_extensions_list or mime_type in supported_binary_mime_types_list:
            supported_type = 'binary'
        else:
            logger.warning(f"Unsupported file type for conversion: {file_path.name} ({file_ext}, {mime_type})")
        
        return {
            'supported_type': supported_type,
            'file_ext': file_ext,
            'mime_type': mime_type
        }

    def create_document_index(
        self, 
        file_path: Path, 
        index_dir: Path, 
        description: str, 
        mime_type: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Create a BM25 index for a single document.
        
        Args:
            file_path: Path to the document
            index_dir: Directory to store the index
            description: Description of the document
            mime_type: Optional MIME type of the file (used when file extension is not available)
            
        Returns:
            Dictionary with index information
        """
        logger.info(f"Processing document: {file_path.name}")
        
        # Validate document type
        validation_result = self.validate_doc_type(file_path, mime_type)
        if validation_result['supported_type'] == 'unsupported':
            logger.warning(f"Cannot create index for unsupported file type: {file_path.name}")
            return {"message": "unsupported file type for indexing"}

        # Convert to text (with page tracking for PDFs)
        text_content, file_type, page_map, total_pages = self.convert_to_text(
            file_path, 
            validation_result['supported_type'],
            validation_result['file_ext'], 
            validation_result['mime_type'],
        )
        
        if not text_content:
            logger.warning(f"No content extracted from {file_path.name}")
            return {"message": "no content extracted from file"}
        
        # Determine position type based on file type
        position_type = 'page'  # default
        if file_type in ['.pptx', '.ppt', 'application/vnd.openxmlformats-officedocument.presentationml.presentation', 'application/vnd.ms-powerpoint']:
            position_type = 'slide'
        elif file_type in ['.txt', '.md', '.csv', '.html', '.htm', '.json', '.xml', 'text/plain', 'text/markdown', 'text/csv', 'text/html', 'application/json', 'application/xml']:
            position_type = 'line'
        
        # Create metadata
        doc_metadata = {
            'filename': file_path.name,
            'filepath': str(file_path),
            'file_type': file_type,
            'size_bytes': file_path.stat().st_size
        }
        
        # Add count with appropriate field name based on position type
        if total_pages:
            if position_type == 'line':
                doc_metadata['total_lines'] = total_pages
            elif position_type == 'slide':
                doc_metadata['total_slides'] = total_pages
            else:
                doc_metadata['total_pages'] = total_pages
        
        # Add description to document metadata - mandatory as positional argument
        doc_metadata['description'] = description
        
        # Chunk the document (with position mapping and type)
        chunks = self.chunker.chunk_text(text_content, doc_metadata, page_map, position_type)
        
        if not chunks:
            logger.warning(f"No chunks created from {file_path.name}")
            return {"message": "no chunks created from file"}
        
        # Extract text from chunks for indexing
        corpus_texts = [chunk['text'] for chunk in chunks]
        
        # Tokenize the corpus (BM25s requires tokenized input)
        # We'll use simple whitespace tokenization with stemming
        corpus_tokens = bm25s.tokenize(
            corpus_texts,
            stemmer=self.stemmer,
            stopwords="en"  # Remove common English stopwords
        )
        
        # Create BM25 index
        retriever = bm25s.BM25()
        retriever.index(corpus_tokens)
        
        # Save the index
        retriever.save(str(index_dir))
        
        # Save metadata and chunks
        metadata = {
            'document': doc_metadata,
            'chunks': chunks,
            'num_chunks': len(chunks),
            'total_chars': len(text_content),
            'index_dir': str(index_dir)
        }
        
        with open(index_dir / "metadata.json", 'w') as f:
            json.dump(metadata, f, indent=2)
        
        # Save corpus for reference
        with open(index_dir / "corpus.pkl", 'wb') as f:
            pickle.dump(corpus_texts, f)
        
        logger.info(f"Created index for {file_path.name} with {len(chunks)} chunks")
        logger.info(f"Index saved to: {index_dir}")
        
        return metadata

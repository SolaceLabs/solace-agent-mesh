"""
File Conversion Service for BM25 Indexing

Converts binary files (PDF, DOCX, PPTX) to text with citation tracking.
Uses BaseArtifactService interface for storage-agnostic operations.

License Compliance:
- Uses pypdf (BSD-3-Clause) instead of PyMuPDF (AGPL v3)
- python-docx (MIT License)
- python-pptx (MIT License)
"""

import asyncio
import logging
from datetime import datetime, timezone
from io import BytesIO
from typing import Tuple, Dict, Any, Optional
from google.adk.artifacts import BaseArtifactService

log = logging.getLogger(__name__)

# Citation configuration for text files
LINES_PER_CITATION_CHUNK = 50  # Group 50 lines per citation entry (reasonable granularity)


def create_line_range_citations(text: str, lines_per_chunk: int = LINES_PER_CITATION_CHUNK) -> dict:
    """
    Create line-range citations for text files (code, config, markdown, etc.).

    Groups lines into chunks for citation purposes. This provides granular
    location information for text files similar to page numbers for PDFs.

    Args:
        text: Full text content
        lines_per_chunk: How many lines per citation entry (default: 50)

    Returns:
        Metadata dict with citation_map

    Example:
        200-line file with lines_per_chunk=50:
        citation_map: [
            {"location": "lines_1_50", "char_start": 0, "char_end": 2500},
            {"location": "lines_51_100", "char_start": 2500, "char_end": 5200},
            {"location": "lines_101_150", "char_start": 5200, "char_end": 7800},
            {"location": "lines_151_200", "char_start": 7800, "char_end": 10000}
        ]
    """
    if not text:
        return {
            "citation_type": "line_range",
            "line_count": 0,
            "char_count": 0,
            "citation_map": [],
            "lines_per_chunk": lines_per_chunk
        }

    lines = text.split('\n')
    citation_map = []
    char_position = 0

    for start_line_idx in range(0, len(lines), lines_per_chunk):
        end_line_idx = min(start_line_idx + lines_per_chunk, len(lines))

        # Get lines for this chunk
        chunk_lines = lines[start_line_idx:end_line_idx]
        chunk_text = '\n'.join(chunk_lines)

        # Calculate character positions
        char_start = char_position
        char_end = char_position + len(chunk_text)

        # Create citation entry (1-indexed for user-friendliness)
        citation_map.append({
            "location": f"lines_{start_line_idx + 1}_{end_line_idx}",
            "char_start": char_start,
            "char_end": char_end
        })

        # Move position past the chunk text AND the newline separator (except for last chunk)
        if end_line_idx < len(lines):
            char_position = char_end + 1  # +1 for newline separator between citation chunks
        else:
            char_position = char_end  # No separator after last chunk

    metadata = {
        "citation_type": "line_range",
        "line_count": len(lines),
        "char_count": len(text),
        "citation_map": citation_map,
        "lines_per_chunk": lines_per_chunk,
        "timestamp": datetime.now(timezone.utc).isoformat()
    }

    log.debug(
        f"Text file citations: {len(lines)} lines → {len(citation_map)} line-range citations "
        f"({lines_per_chunk} lines per citation)"
    )

    return metadata


def convert_pdf_to_text(pdf_bytes: bytes) -> Tuple[str, dict]:
    """
    Extract text from PDF using pypdf with page-level citation tracking.

    IMPORTANT: Page numbers are PHYSICAL/SEQUENTIAL (1, 2, 3, ...),
    NOT the document's internal page numbers (which may use i, ii, iii, 1, 2, etc.).

    Example:
      Physical page 1 → "page_1" (may be titled "Preface" or "i" in document)
      Physical page 2 → "page_2" (may be "ii" in document)
      Physical page 10 → "page_10" (may be "Page 1" in document)

    This ensures consistent, unambiguous page references.

    Args:
        pdf_bytes: PDF file content as bytes

    Returns:
        Tuple of (extracted_text, conversion_metadata)

        conversion_metadata includes:
        - converter: "pypdf"
        - page_count: number of pages
        - char_count: total characters extracted
        - citation_map: list of {"location": "page_N", "char_start": X, "char_end": Y}
          where N is the physical/sequential page number (1, 2, 3, ...)
    """
    try:
        from pypdf import PdfReader

        reader = PdfReader(BytesIO(pdf_bytes))

        full_text = []
        citation_map = []
        char_position = 0

        for page_num, page in enumerate(reader.pages, start=1):
            # Extract text from page
            page_text = page.extract_text()

            if page_text:
                # Record citation info
                char_start = char_position
                char_end = char_position + len(page_text)

                # Use "physical_page_N" to clarify this is sequential page number,
                # not the document's internal page number (which may be i, ii, iii, etc.)
                citation_map.append({
                    "location": f"physical_page_{page_num}",
                    "char_start": char_start,
                    "char_end": char_end
                })

                full_text.append(page_text)

                # CRITICAL: Account for newline that will be added by "\n".join()
                # Move position past the page text AND the newline character
                char_position = char_end + 1  # +1 for the newline between pages

        # Join with newlines (this is why we added +1 above)
        extracted_text = "\n".join(full_text)

        # IMPORTANT: Fix the last citation's char_end (no newline after last page)
        if citation_map:
            # The last page doesn't have a trailing newline, so subtract 1
            citation_map[-1]["char_end"] = len(extracted_text)

        metadata = {
            "converter": "pypdf",
            "page_count": len(reader.pages),
            "char_count": len(extracted_text),
            "citation_type": "page",
            "citation_map": citation_map,
            "timestamp": datetime.now(timezone.utc).isoformat()
        }

        log.info(
            f"PDF conversion: {len(reader.pages)} pages, {len(extracted_text)} chars, "
            f"{len(citation_map)} citations"
        )

        return extracted_text, metadata

    except Exception as e:
        log.error(f"PDF conversion failed: {e}")
        raise ValueError(f"Failed to convert PDF: {e}") from e


def convert_docx_to_text(docx_bytes: bytes) -> Tuple[str, dict]:
    """
    Extract text from DOCX using python-docx with paragraph-level citation tracking.

    IMPORTANT: Paragraph numbers are PHYSICAL/SEQUENTIAL (1, 2, 3, ...),
    representing the order of paragraphs in the document.

    This ensures consistent, unambiguous paragraph references.

    Args:
        docx_bytes: DOCX file content as bytes

    Returns:
        Tuple of (extracted_text, conversion_metadata)

        conversion_metadata includes:
        - converter: "python-docx"
        - paragraph_count: number of paragraphs
        - char_count: total characters extracted
        - citation_map: list of {"location": "physical_paragraph_N", "char_start": X, "char_end": Y}
          where N is the physical/sequential paragraph number (1, 2, 3, ...)
    """
    try:
        from docx import Document

        doc = Document(BytesIO(docx_bytes))

        full_text = []
        citation_map = []
        char_position = 0

        for para_num, paragraph in enumerate(doc.paragraphs, start=1):
            para_text = paragraph.text

            if para_text.strip():  # Only include non-empty paragraphs
                # Record citation info
                char_start = char_position
                char_end = char_position + len(para_text)

                # Use "physical_paragraph_N" for consistency with physical_page_N
                citation_map.append({
                    "location": f"physical_paragraph_{para_num}",
                    "char_start": char_start,
                    "char_end": char_end
                })

                full_text.append(para_text)

                # CRITICAL: Account for newline that will be added by "\n".join()
                char_position = char_end + 1  # +1 for the newline between paragraphs

        # Join with newlines (this is why we added +1 above)
        extracted_text = "\n".join(full_text)

        # IMPORTANT: Fix the last citation's char_end (no newline after last paragraph)
        if citation_map:
            citation_map[-1]["char_end"] = len(extracted_text)

        metadata = {
            "converter": "python-docx",
            "paragraph_count": len([p for p in doc.paragraphs if p.text.strip()]),
            "char_count": len(extracted_text),
            "citation_type": "paragraph",
            "citation_map": citation_map,
            "timestamp": datetime.now(timezone.utc).isoformat()
        }

        log.info(
            f"DOCX conversion: {len(doc.paragraphs)} paragraphs, {len(extracted_text)} chars, "
            f"{len(citation_map)} citations"
        )

        return extracted_text, metadata

    except Exception as e:
        log.error(f"DOCX conversion failed: {e}")
        raise ValueError(f"Failed to convert DOCX: {e}") from e


def convert_pptx_to_text(pptx_bytes: bytes) -> Tuple[str, dict]:
    """
    Extract text from PPTX using python-pptx with slide-level citation tracking.

    IMPORTANT: Slide numbers are PHYSICAL/SEQUENTIAL (1, 2, 3, ...),
    representing the physical order of slides in the presentation,
    NOT any custom slide numbers displayed in the presentation.

    Example:
      Physical slide 1 → "physical_slide_1" (may show "Title" or no number)
      Physical slide 2 → "physical_slide_2" (may show "Slide 1" internally)
      Physical slide 10 → "physical_slide_10" (regardless of custom numbering)

    This ensures consistent, unambiguous slide references, especially when
    slides are hidden or custom numbering is used.

    Args:
        pptx_bytes: PPTX file content as bytes

    Returns:
        Tuple of (extracted_text, conversion_metadata)

        conversion_metadata includes:
        - converter: "python-pptx"
        - slide_count: number of slides
        - char_count: total characters extracted
        - citation_map: list of {"location": "physical_slide_N", "char_start": X, "char_end": Y}
          where N is the physical/sequential slide number (1, 2, 3, ...)
    """
    try:
        from pptx import Presentation

        prs = Presentation(BytesIO(pptx_bytes))

        full_text = []
        citation_map = []
        char_position = 0

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text_parts = []

            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text_parts.append(shape.text)

            if slide_text_parts:
                slide_text = "\n".join(slide_text_parts)

                # Record citation info
                char_start = char_position
                char_end = char_position + len(slide_text)

                # Use "physical_slide_N" to clarify this is sequential slide number,
                # not any custom slide numbering that may be displayed
                citation_map.append({
                    "location": f"physical_slide_{slide_num}",
                    "char_start": char_start,
                    "char_end": char_end
                })

                full_text.append(slide_text)

                # CRITICAL: Account for double newline that will be added by "\n\n".join()
                # Move position past the slide text AND the two newline characters
                char_position = char_end + 2  # +2 for the double newline between slides

        # Join with double newlines (this is why we added +2 above)
        extracted_text = "\n\n".join(full_text)

        # IMPORTANT: Fix the last citation's char_end (no double newline after last slide)
        if citation_map:
            citation_map[-1]["char_end"] = len(extracted_text)

        metadata = {
            "converter": "python-pptx",
            "slide_count": len(prs.slides),
            "char_count": len(extracted_text),
            "citation_type": "slide",
            "citation_map": citation_map,
            "timestamp": datetime.now(timezone.utc).isoformat()
        }

        log.info(
            f"PPTX conversion: {len(prs.slides)} slides, {len(extracted_text)} chars, "
            f"{len(citation_map)} citations"
        )

        return extracted_text, metadata

    except Exception as e:
        log.error(f"PPTX conversion failed: {e}")
        raise ValueError(f"Failed to convert PPTX: {e}") from e


def should_convert_file(mime_type: str, filename: str) -> bool:
    """
    Check if file should be converted based on MIME type.

    Returns True for:
    - application/pdf (PDF files)
    - application/vnd.openxmlformats-officedocument.wordprocessingml.document (DOCX)
    - application/vnd.openxmlformats-officedocument.presentationml.presentation (PPTX)

    Args:
        mime_type: The MIME type of the file
        filename: The filename (for future extension-based fallback if needed)

    Returns:
        bool: True if file should be converted, False otherwise
    """
    if not mime_type:
        return False

    convertible_types = [
        "application/pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",  # DOCX
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"  # PPTX
    ]

    return mime_type in convertible_types


async def convert_and_save_artifact(
    artifact_service: BaseArtifactService,
    app_name: str,
    user_id: str,
    session_id: str,
    source_filename: str,
    source_version: int,
    mime_type: str
) -> dict:
    """
    Load source artifact, convert to text, and save as new artifact.

    STORAGE-AGNOSTIC: Uses BaseArtifactService interface.
    Works with S3, GCS, or filesystem backends transparently.

    Args:
        artifact_service: The artifact service instance
        app_name: Application name
        user_id: User ID
        session_id: Session ID (e.g., "project-{project_id}")
        source_filename: Original filename (e.g., "report.pdf")
        source_version: Version of the source file
        mime_type: MIME type of the source file

    Returns:
        Result dict with status:
        - Success: {"status": "success", "data_version": int, ...}
        - Error: {"status": "error", "error": "descriptive error message"}
        - Not convertible: {"status": "skipped", "reason": "MIME type not supported"}
    """
    log_prefix = f"[FileConverter:{source_filename}:v{source_version}]"

    # Check if file should be converted
    if not should_convert_file(mime_type, source_filename):
        log.debug(f"{log_prefix} MIME type {mime_type} not convertible, skipping")
        return {
            "status": "skipped",
            "reason": f"MIME type {mime_type} is not supported for conversion"
        }

    log.info(f"{log_prefix} Starting conversion for {mime_type}")

    try:
        # Import artifact helpers
        from ....agent.utils.artifact_helpers import (
            load_artifact_content_or_metadata,
            save_artifact_with_metadata
        )

        # 1. Load source binary artifact using artifact service (storage-agnostic)
        source_result = await load_artifact_content_or_metadata(
            artifact_service=artifact_service,
            app_name=app_name,
            user_id=user_id,
            session_id=session_id,
            filename=source_filename,
            version=source_version,
            return_raw_bytes=True
        )

        if source_result.get("status") != "success":
            error_msg = f"Failed to load source file '{source_filename}': {source_result.get('message', 'Unknown error')}"
            log.error(f"{log_prefix} {error_msg}")
            return {"status": "error", "error": error_msg}

        source_bytes = source_result.get("raw_bytes")
        if not source_bytes:
            error_msg = f"Source file '{source_filename}' has no content"
            log.error(f"{log_prefix} {error_msg}")
            return {"status": "error", "error": error_msg}

        # 2. Detect type and call appropriate converter (in-memory processing)
        # Run conversion in thread pool to avoid blocking event loop and allow SSE events to be sent
        try:
            loop = asyncio.get_running_loop()
            if mime_type == "application/pdf":
                text, conversion_metadata = await loop.run_in_executor(
                    None, convert_pdf_to_text, source_bytes
                )
            elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                text, conversion_metadata = await loop.run_in_executor(
                    None, convert_docx_to_text, source_bytes
                )
            elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                text, conversion_metadata = await loop.run_in_executor(
                    None, convert_pptx_to_text, source_bytes
                )
            else:
                error_msg = f"Unsupported MIME type for conversion: {mime_type}"
                log.warning(f"{log_prefix} {error_msg}")
                return {"status": "error", "error": error_msg}
        except ValueError as e:
            # Conversion-specific errors (corrupted files, password-protected, etc.)
            file_type = mime_type.split('/')[-1].upper().replace('APPLICATION/VND.OPENXMLFORMATS-OFFICEDOCUMENT.', '')
            error_msg = f"Failed to convert {file_type} '{source_filename}': {str(e)}"
            log.error(f"{log_prefix} Conversion error: {e}")
            return {"status": "error", "error": error_msg}
        except Exception as e:
            # Unexpected errors during conversion
            error_msg = f"Unexpected error converting '{source_filename}': {str(e)}"
            log.exception(f"{log_prefix} {error_msg}")
            return {"status": "error", "error": error_msg}

        # 3. Generate converted filename
        converted_filename = f"{source_filename}.converted.txt"

        # 4. Save using artifact service (works with S3/GCS/filesystem)
        # CRITICAL: source_version matches converted version (both created in same operation)
        try:
            result = await save_artifact_with_metadata(
                artifact_service=artifact_service,
                app_name=app_name,
                user_id=user_id,
                session_id=session_id,
                filename=converted_filename,
                content_bytes=text.encode('utf-8'),
                mime_type="text/plain",
                metadata_dict={
                    "source": "conversion",
                    "conversion": {
                        "source_file": source_filename,
                        "source_version": source_version,  # Matches converted file version
                        **conversion_metadata
                    }
                },
                timestamp=datetime.now(timezone.utc)
            )

            if result.get("status") != "success":
                error_msg = f"Failed to save converted file '{converted_filename}': {result.get('status_message', 'Unknown error')}"
                log.error(f"{log_prefix} {error_msg}")
                return {"status": "error", "error": error_msg}

            log.info(
                f"{log_prefix} Converted to {converted_filename} v{result.get('data_version')} "
                f"({conversion_metadata.get('char_count')} chars, "
                f"{len(conversion_metadata.get('citation_map', []))} citations)"
            )

            return result

        except Exception as e:
            # Errors during save operation
            error_msg = f"Failed to save converted file '{converted_filename}': {str(e)}"
            log.exception(f"{log_prefix} {error_msg}")
            return {"status": "error", "error": error_msg}

    except Exception as e:
        # Catch-all for any errors not caught by inner try/except blocks
        error_msg = f"Unexpected error in conversion pipeline for '{source_filename}': {str(e)}"
        log.exception(f"{log_prefix} {error_msg}")
        return {"status": "error", "error": error_msg}

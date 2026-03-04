"""Shared pytest fixtures for BM25 service tests."""

import pytest
from unittest.mock import MagicMock, AsyncMock
from io import BytesIO


@pytest.fixture
def mock_artifact_service():
    """Create a mock artifact service for testing."""
    service = AsyncMock()
    service.load_artifact = AsyncMock()
    service.save_artifact = AsyncMock()
    service.delete_artifact = AsyncMock()
    service.list_artifacts = AsyncMock()
    return service


@pytest.fixture
def mock_sse_manager():
    """Create a mock SSE manager for task service tests."""
    manager = MagicMock()
    manager.send_event = AsyncMock()
    return manager


@pytest.fixture
def mock_project_service():
    """Create a mock project service for task service tests."""
    service = MagicMock()
    service.artifact_service = AsyncMock()
    service.app_name = "test-app"
    return service


@pytest.fixture
def sample_pdf_bytes():
    """Create sample PDF bytes for testing."""
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter

        pdf_buffer = BytesIO()
        c = canvas.Canvas(pdf_buffer, pagesize=letter)
        c.drawString(100, 750, "Test PDF content")
        c.showPage()
        c.save()

        return pdf_buffer.getvalue()
    except ImportError:
        # If reportlab not available, return minimal valid PDF
        return b"%PDF-1.4\n1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj 2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj 3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R/Resources<<>>>>endobj\nxref\n0 4\n0000000000 65535 f\n0000000009 00000 n\n0000000056 00000 n\n0000000115 00000 n\ntrailer<</Size 4/Root 1 0 R>>\nstartxref\n203\n%%EOF"


@pytest.fixture
def sample_docx_bytes():
    """Create sample DOCX bytes for testing."""
    try:
        from docx import Document

        doc = Document()
        doc.add_paragraph("Test DOCX paragraph 1")
        doc.add_paragraph("Test DOCX paragraph 2")

        docx_buffer = BytesIO()
        doc.save(docx_buffer)
        docx_buffer.seek(0)

        return docx_buffer.getvalue()
    except ImportError:
        pytest.skip("python-docx not available")


@pytest.fixture
def sample_pptx_bytes():
    """Create sample PPTX bytes for testing."""
    try:
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Test Slide 1"

        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        return pptx_buffer.getvalue()
    except ImportError:
        pytest.skip("python-pptx not available")


@pytest.fixture
def mock_project():
    """Create a mock project entity for testing."""
    project = MagicMock()
    project.id = "test-project-123"
    project.user_id = "test-user-456"
    project.name = "Test Project"
    return project
